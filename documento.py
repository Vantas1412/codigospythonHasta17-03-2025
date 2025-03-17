from pptx import Presentation
from pptx.util import Inches

def crear_presentacion(contenido_intro, contenido_desarrollo, imagenes, nombre_archivo="presentacion.pptx"):
    prs = Presentation()

  
    for img_path in imagenes:
        slide = prs.slides.add_slide(prs.slide_layouts[5])  
        left = Inches(1)
        top = Inches(1)
        width = Inches(8)
        slide.shapes.add_picture(img_path, left, top, width=width)
    
    prs.save(nombre_archivo)
    print(f"Presentación '{nombre_archivo}' creada exitosamente.")


imagen1="C:/Users/Usuario/Pictures/Screenshots/MenfiTarea/pagina1.png"
imagen2="C:/Users/Usuario/Pictures/Screenshots/MenfiTarea/pagina2.png"
imagen3="C:/Users/Usuario/Pictures/Screenshots/MenfiTarea/pagina3.png"
imagen4="C:/Users/Usuario/Pictures/Screenshots/MenfiTarea/pagina4.png"
imagen5="C:/Users/Usuario/Pictures/Screenshots/MenfiTarea/pagina5.png"
imagenes = [imagen1, imagen2, imagen3, imagen4, imagen5]
crear_presentacion("Contenido introductorio", "Desarrollo del tema principal", imagenes)
